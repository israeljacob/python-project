from pptx import Presentation


def extract_text_from_text_frame(shape):
    """
    Extracts text from a text frame shape.
    Args:
        shape: The text frame shape object.
    Returns:
        str: The extracted text from the text frame.
    """
    return " ".join(run.text for paragraph in shape.text_frame.paragraphs for run in paragraph.runs) + " "


def extract_text_from_table(shape):
    """
    Extracts text from a table shape.
    Args:
        shape: The table shape object.
    Returns:
        str: The extracted text from the table.
    """
    extracted_text = ""
    table = shape.table
    for row in table.rows:
        for cell in row.cells:
            extracted_text += " ".join(
                run.text for paragraph in cell.text_frame.paragraphs for run in paragraph.runs) + " "
    return extracted_text


def extract_text_from_slide(slide):
    """
    Extracts text from a slide.
    Args:
        slide: The slide object.
    Returns:
        str: The extracted text from the slide.
    """
    extracted_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            extracted_text += extract_text_from_text_frame(shape)
        elif shape.has_table:
            extracted_text += extract_text_from_table(shape)
    return extracted_text


def extract_text_from_pptx(path_file):
    """
    Extracts text from a .pptx file using the python-pptx package.
    Args:
        path_file (str): The path to the .pptx file.
    Returns:
        dict: A dictionary where the slide number is the key and the extracted text is the value.
    """
    presentation = Presentation(path_file)
    slide_texts = {}
    for i, slide in enumerate(presentation.slides, start=1):
        slide_texts[i] = extract_text_from_slide(slide)
    return slide_texts
